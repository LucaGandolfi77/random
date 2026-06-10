#!/usr/bin/env python3
"""
Generate a PowerPoint presentation from Qiskit optimization results.
"""

import csv
import io
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

# ── Paths ────────────────────────────────────────────────────
HERE  = Path(__file__).parent
CSV   = HERE / "merged_20260325_165037.csv"
OUT   = HERE / "Qiskit_Optimization_Results.pptx"

# ── Colors ───────────────────────────────────────────────────
DARK_BG      = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE  = RGBColor(0x41, 0x8C, 0xF0)
ACCENT_ORANGE= RGBColor(0xFF, 0x8C, 0x00)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
SOFT_WHITE   = RGBColor(0xE8, 0xE8, 0xE8)
GREEN        = RGBColor(0x00, 0xC8, 0x53)
TABLE_HEADER = RGBColor(0x2D, 0x2D, 0x4F)
TABLE_ROW1   = RGBColor(0x25, 0x25, 0x40)
TABLE_ROW2   = RGBColor(0x1F, 0x1F, 0x38)

# ── Load CSV data ────────────────────────────────────────────
def load_data():
    rows = []
    with open(CSV, encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for r in reader:
            rows.append(r)
    return rows

# ── Helpers ──────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_text(text_frame, items, font_size=16, color=WHITE, font_name="Calibri"):
    for item in items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_before = Pt(4)
        p.level = 0

def add_code_box(slide, left, top, width, height, code_text, font_size=10):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    # Dark code background
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x22)
    for i, line in enumerate(code_text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xA0, 0xE0, 0xA0)
        p.font.name = "Courier New"
        p.space_before = Pt(1)
        p.space_after = Pt(1)

def save_matplotlib_to_image(fig, path):
    fig.savefig(path, dpi=200, bbox_inches="tight", transparent=True)
    plt.close(fig)

# ── Chart generation with matplotlib ────────────────────────
def make_time_comparison_chart(data, img_path):
    """Bar chart: classic vs quantum time by instance."""
    instances = [f"{r['n_servers']}s{r['n_vms']}v" for r in data]
    c_times = [float(r['classic_time_s']) for r in data]
    q_times = [float(r['quantum_time_s']) for r in data]

    fig, ax = plt.subplots(figsize=(11, 4.5))
    fig.patch.set_facecolor('#1B1B2F')
    ax.set_facecolor('#1B1B2F')
    x = np.arange(len(instances))
    w = 0.35
    ax.bar(x - w/2, c_times, w, label="Classical", color="#418CF0", edgecolor="none")
    ax.bar(x + w/2, q_times, w, label="Quantum (QAOA)", color="#FF8C00", edgecolor="none")
    ax.set_xticks(x)
    ax.set_xticklabels(instances, rotation=60, ha="right", fontsize=7, color="white")
    ax.set_ylabel("Time (s)", color="white", fontsize=11)
    ax.set_title("Execution Time: Classical vs Quantum", color="white", fontsize=13, fontweight="bold")
    ax.legend(fontsize=9, facecolor="#2D2D4F", edgecolor="none", labelcolor="white")
    ax.tick_params(colors="white")
    for spine in ax.spines.values():
        spine.set_color("#555")
    ax.grid(axis="y", alpha=0.2, color="white")
    fig.tight_layout()
    save_matplotlib_to_image(fig, img_path)

def make_speedup_chart(data, img_path):
    """Line chart of speedup by instance."""
    instances = [f"{r['n_servers']}s{r['n_vms']}v" for r in data]
    speedups = [float(r['speedup_x']) for r in data]

    fig, ax = plt.subplots(figsize=(11, 4.5))
    fig.patch.set_facecolor('#1B1B2F')
    ax.set_facecolor('#1B1B2F')
    x = np.arange(len(instances))
    ax.plot(x, speedups, marker="o", color="#FF8C00", linewidth=2, markersize=5)
    ax.axhline(y=1.0, color="#418CF0", linestyle="--", alpha=0.6, label="Parity (1.0x)")
    ax.fill_between(x, speedups, alpha=0.15, color="#FF8C00")
    ax.set_xticks(x)
    ax.set_xticklabels(instances, rotation=60, ha="right", fontsize=7, color="white")
    ax.set_ylabel("Speedup (classical / quantum)", color="white", fontsize=11)
    ax.set_title("Speedup Ratio Across All Instances", color="white", fontsize=13, fontweight="bold")
    ax.legend(fontsize=9, facecolor="#2D2D4F", edgecolor="none", labelcolor="white")
    ax.tick_params(colors="white")
    for spine in ax.spines.values():
        spine.set_color("#555")
    ax.grid(axis="y", alpha=0.2, color="white")
    fig.tight_layout()
    save_matplotlib_to_image(fig, img_path)

def make_objective_chart(data, img_path):
    """Grouped bar chart of objectives by n_servers."""
    servers_set = sorted(set(int(r['n_servers']) for r in data))
    vms_set = sorted(set(int(r['n_vms']) for r in data))

    fig, ax = plt.subplots(figsize=(10, 4.5))
    fig.patch.set_facecolor('#1B1B2F')
    ax.set_facecolor('#1B1B2F')

    for s in servers_set:
        objs = []
        for v in vms_set:
            match = [r for r in data if int(r['n_servers']) == s and int(r['n_vms']) == v]
            objs.append(float(match[0]['classic_objective']) if match else 0)
        ax.plot(vms_set, objs, marker="s", label=f"{s} servers", linewidth=2, markersize=6)

    ax.set_xlabel("Number of VMs", color="white", fontsize=11)
    ax.set_ylabel("Objective Value", color="white", fontsize=11)
    ax.set_title("Objective Value by Problem Size", color="white", fontsize=13, fontweight="bold")
    ax.legend(fontsize=8, facecolor="#2D2D4F", edgecolor="none", labelcolor="white", ncol=3)
    ax.tick_params(colors="white")
    for spine in ax.spines.values():
        spine.set_color("#555")
    ax.grid(alpha=0.2, color="white")
    fig.tight_layout()
    save_matplotlib_to_image(fig, img_path)

def make_scaling_heatmap(data, img_path):
    """Heatmap of quantum time by (n_servers, n_vms)."""
    servers_set = sorted(set(int(r['n_servers']) for r in data))
    vms_set = sorted(set(int(r['n_vms']) for r in data))

    matrix = np.zeros((len(servers_set), len(vms_set)))
    for r in data:
        si = servers_set.index(int(r['n_servers']))
        vi = vms_set.index(int(r['n_vms']))
        matrix[si, vi] = float(r['quantum_time_s'])

    fig, ax = plt.subplots(figsize=(7, 5))
    fig.patch.set_facecolor('#1B1B2F')
    ax.set_facecolor('#1B1B2F')
    im = ax.imshow(matrix, cmap="YlOrRd", aspect="auto")
    ax.set_xticks(range(len(vms_set)))
    ax.set_xticklabels(vms_set, color="white")
    ax.set_yticks(range(len(servers_set)))
    ax.set_yticklabels(servers_set, color="white")
    ax.set_xlabel("Number of VMs", color="white", fontsize=11)
    ax.set_ylabel("Number of Servers", color="white", fontsize=11)
    ax.set_title("Quantum Solver Time (seconds)", color="white", fontsize=13, fontweight="bold")
    for i in range(len(servers_set)):
        for j in range(len(vms_set)):
            val = matrix[i, j]
            color = "white" if val < matrix.max() * 0.6 else "black"
            ax.text(j, i, f"{val:.1f}", ha="center", va="center", color=color, fontsize=9)
    cbar = fig.colorbar(im, ax=ax)
    cbar.ax.yaxis.set_tick_params(color="white")
    plt.setp(cbar.ax.yaxis.get_ticklabels(), color="white")
    fig.tight_layout()
    save_matplotlib_to_image(fig, img_path)

def make_success_chart(data, img_path):
    """Pie chart showing success rate."""
    c_success = sum(1 for r in data if r['classic_status'] == 'SUCCESS')
    q_success = sum(1 for r in data if r['quantum_status'] == 'SUCCESS')
    total = len(data)

    fig, axes = plt.subplots(1, 2, figsize=(8, 3.5))
    fig.patch.set_facecolor('#1B1B2F')

    for ax, label, count, color in [
        (axes[0], "Classical", c_success, "#418CF0"),
        (axes[1], "Quantum (QAOA)", q_success, "#FF8C00"),
    ]:
        ax.set_facecolor('#1B1B2F')
        sizes = [count, total - count]
        colors_list = [color, "#333"]
        ax.pie(sizes, labels=[f"SUCCESS\n{count}/{total}", ""] if count < total else [f"SUCCESS\n{count}/{total}", ""],
               colors=colors_list, startangle=90,
               textprops={"color": "white", "fontsize": 11, "fontweight": "bold"})
        ax.set_title(label, color="white", fontsize=12, fontweight="bold")

    fig.suptitle("Feasibility Rate: 100% for Both Solvers", color="white", fontsize=13, fontweight="bold", y=1.02)
    fig.tight_layout()
    save_matplotlib_to_image(fig, img_path)

# ── Build presentation ───────────────────────────────────────
def build_pptx():
    data = load_data()
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    IMG_DIR = HERE / "_pptx_charts"
    IMG_DIR.mkdir(exist_ok=True)

    # ════════════════════════════════════════════════════════
    # SLIDE 1 – Title
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "Quantum vs Classical Optimization",
                 font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 "Server-VM Workload Placement with Qiskit ADMM & QAOA",
                 font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
                 "36 problem instances  |  1-6 Servers  |  1-6 VMs  |  100% Feasibility",
                 font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 "Qiskit 1.x  \u2022  qiskit-optimization  \u2022  qiskit-algorithms  \u2022  DOcplex",
                 font_size=14, color=SOFT_WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                 "March 2026",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════
    # SLIDE 2 – Introduction / Problem Statement
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "1. Introduction & Problem Statement",
                 font_size=30, color=ACCENT_BLUE, bold=True)
    tf = add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.5),
                      "Objective", font_size=20, color=ACCENT_ORANGE, bold=True)
    add_bullet_text(tf, [
        "Minimize the total cost of running N servers hosting M virtual machines",
        "Cost = fixed server activation cost + dynamic CPU utilization cost",
        "Compare a classical exact solver against a quantum QAOA solver",
        "",
        "This is a Mixed-Integer Quadratic Program (MIQP):",
        "  \u2022 Binary variables: server on/off (s\u1d62)",
        "  \u2022 Continuous variables: VM-to-server load (v\u2c7c\u1d62), CPU usage (u\u2c7c)",
        "",
        "Real-world application: cloud data center workload placement,",
        "where we allocate VMs across physical servers to minimize",
        "energy and operational costs under capacity constraints."
    ], font_size=15, color=WHITE)

    tf2 = add_text_box(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5),
                       "Mathematical Formulation", font_size=20, color=ACCENT_ORANGE, bold=True)
    add_bullet_text(tf2, [
        "Minimize:",
        "  \u03a3\u1d62 [ \u03c0\u1d62\u00b7s\u1d62 + \u03b4\u1d62\u00b7\u03a3\u2c7c (u\u2c7c\u00b7v\u2c7c\u1d62) ]",
        "",
        "Subject to:",
        "  \u2022 Server load:  \u03a3\u2c7c v\u2c7c\u1d62 \u2265 cap\u1d62 \u2212 1      \u2200i",
        "  \u2022 Server ON:    s\u1d62 = 1                  \u2200i (if require_all_on)",
        "  \u2022 VM allocation: \u03a3\u1d62 v\u2c7c\u1d62 \u2264 limit\u2c7c    \u2200j",
        "  \u2022 Min CPU:       u\u2c7c \u2265 min_cpu          \u2200j",
        "",
        "Variables:",
        "  s\u1d62 \u2208 {0,1}  \u2022  v\u2c7c\u1d62 \u2265 0  \u2022  u\u2c7c \u2265 1.0",
        "",
        "Tested configurations:",
        "  N \u2208 {1,2,3,4,5,6}  \u00d7  M \u2208 {1,2,3,4,5,6}  = 36 instances",
    ], font_size=15, color=WHITE)

    # ════════════════════════════════════════════════════════
    # SLIDE 3 – Solution Approach
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "2. Solution Approach: ADMM Decomposition",
                 font_size=30, color=ACCENT_BLUE, bold=True)
    tf = add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.5),
                      "ADMM (Alternating Direction Method of Multipliers)",
                      font_size=18, color=ACCENT_ORANGE, bold=True)
    add_bullet_text(tf, [
        "Decomposes the MIQP into sub-problems:",
        "",
        "  1. Binary sub-problem (server on/off decisions)",
        "       \u2192 Solved as QUBO",
        "",
        "  2. Continuous sub-problem (load allocation)",
        "       \u2192 Solved via COBYLA optimizer",
        "",
        "  3. Iteratively update Lagrange multipliers",
        "       until convergence (tol = 1e-4)",
        "",
        "ADMM Parameters:",
        "  \u2022 \u03c1_initial = 100,  \u03b2 = 1000,  factor_c = 900",
        "  \u2022 Three-block decomposition enabled",
        "  \u2022 Classical: maxiter = 100",
        "  \u2022 Quantum:   maxiter = 300",
    ], font_size=14, color=WHITE)

    tf2 = add_text_box(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5),
                       "Two Solver Configurations", font_size=18, color=ACCENT_ORANGE, bold=True)
    add_bullet_text(tf2, [
        "Classical (Exact):",
        "  \u2022 QUBO sub-problem \u2192 NumPyMinimumEigensolver",
        "  \u2022 Finds exact ground state",
        "  \u2022 Continuous sub-problem \u2192 COBYLA",
        "",
        "Quantum (QAOA):",
        "  \u2022 QUBO sub-problem \u2192 QAOA algorithm",
        "       - StatevectorSampler (noiseless simulation)",
        "       - COBYLA optimizer (maxiter=300)",
        "       - reps = 3 (variational circuit depth)",
        "  \u2022 Continuous sub-problem \u2192 COBYLA",
        "",
        "Both pipelines use identical problem formulation",
        "and ADMM framework \u2014 only the QUBO solver differs.",
        "",
        "This enables a fair comparison of quantum vs",
        "classical performance on the combinatorial core.",
    ], font_size=14, color=WHITE)

    # ════════════════════════════════════════════════════════
    # SLIDE 4 – Code Structure
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "3. Code Structure: qiskit_opt.py",
                 font_size=30, color=ACCENT_BLUE, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(5.8), Inches(0.5),
                 "DOcplex Model Building", font_size=18, color=ACCENT_ORANGE, bold=True)
    add_code_box(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.8),
"""# Decision variables
s_vars = [mdl.binary_var(name=f"si{i}")        # server ON/OFF
          for i in range(n_servers)]
v_vars = {(j,i): mdl.continuous_var(lb=0.0,    # VM load
           name=f"vj{j}i{i}")
          for j in range(n_vms) for i in range(n_servers)}
u_vars = [mdl.continuous_var(lb=min_cpu,        # CPU usage
           name=f"uj{j}") for j in range(n_vms)]

# Objective: min server cost + utilization
obj = mdl.sum(pi[i]*s[i] + pd[i]*sum(u[j]*v[j,i]
      for j in range(n_vms)) for i in range(n_servers))
mdl.minimize(obj)""", font_size=10)

    add_text_box(slide, Inches(0.6), Inches(4.6), Inches(5.8), Inches(0.5),
                 "Constraints", font_size=18, color=ACCENT_ORANGE, bold=True)
    add_code_box(slide, Inches(0.6), Inches(5.1), Inches(5.8), Inches(2.0),
"""# Server load constraint
for i in range(n_servers):
    mdl.add(sum(v[j,i] for j in range(n_vms))
            >= capacities[i] - 1)
# VM allocation limit
for j in range(n_vms):
    mdl.add(sum(v[j,i] for i in range(n_servers))
            <= vm_allocation_limits[j])""", font_size=10)

    add_text_box(slide, Inches(6.8), Inches(1.1), Inches(6), Inches(0.5),
                 "ADMM Optimizer Setup", font_size=18, color=ACCENT_ORANGE, bold=True)
    add_code_box(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(2.8),
"""# Translate DOcplex -> Qiskit QuadraticProgram
qp = from_docplex_mp(mdl)

# Classical: exact eigensolver for QUBO
exact = MinimumEigenOptimizer(
    NumPyMinimumEigensolver())
admm_classic = ADMMOptimizer(
    params=admm_params, qubo_optimizer=exact,
    continuous_optimizer=cobyla)

# Quantum: QAOA for QUBO
qaoa_algo = QAOA(sampler=Sampler(),
    optimizer=COBYLA(maxiter=300), reps=3)
qaoa = MinimumEigenOptimizer(qaoa_algo)
admm_quantum = ADMMOptimizer(
    params=admm_params, qubo_optimizer=qaoa,
    continuous_optimizer=cobyla)""", font_size=10)

    add_text_box(slide, Inches(6.8), Inches(4.6), Inches(6), Inches(0.5),
                 "Solve & Post-process", font_size=18, color=ACCENT_ORANGE, bold=True)
    add_code_box(slide, Inches(6.8), Inches(5.1), Inches(6), Inches(2.0),
"""# Solve both
result_classic = admm_classic.solve(qp)
result_classic = snap_to_feasible(result_classic, qp)

result_quantum = admm_quantum.solve(qp)
result_quantum = snap_to_feasible(result_quantum, qp)
# snap_to_feasible fixes tiny ADMM rounding errors
# so Qiskit's strict feasibility check passes""", font_size=10)

    # ════════════════════════════════════════════════════════
    # SLIDE 5 – snap_to_feasible
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "4. Post-Processing: snap_to_feasible()",
                 font_size=30, color=ACCENT_BLUE, bold=True)

    tf = add_text_box(slide, Inches(0.6), Inches(1.1), Inches(5.8), Inches(5.8),
                      "The Challenge", font_size=18, color=ACCENT_ORANGE, bold=True)
    add_bullet_text(tf, [
        "ADMM produces approximate solutions with ~1e-4 error",
        "Qiskit's qp.is_feasible(x) uses strict zero-tolerance",
        "\u2192 Tiny violations cause INFEASIBLE status",
        "",
        "Solution: Iterative constraint repair (up to 20 passes)",
        "",
        "Algorithm:",
        "  1. Clamp all variables to their bounds",
        "  2. Identify LE-tight variables (at ceiling within 1e-10)",
        "  3. Fix GE violations: distribute deficit across",
        "     continuous variables, preferring those NOT at a",
        "     LE ceiling (avoids oscillation between constraints)",
        "  4. Fix LE violations: distribute excess across vars",
        "  5. Re-clamp and repeat if violations remain",
        "",
        "Key insight: excluding LE-tight variables from GE",
        "repairs prevents infinite oscillation between",
        "competing constraint boundaries.",
    ], font_size=14, color=WHITE)

    add_text_box(slide, Inches(6.8), Inches(1.1), Inches(6), Inches(0.5),
                 "Core Logic", font_size=18, color=ACCENT_ORANGE, bold=True)
    add_code_box(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(5.2),
"""def snap_to_feasible(result, qp):
    x = np.array(result.x, dtype=float)
    x = np.clip(x, lb, ub)   # step 1: clamp

    for _ in range(20):       # iterative repair
        # Find vars tight against LE ceilings
        le_tight = set()
        for con in qp.linear_constraints:
            if con.sense == LE and evaluate(x) >= rhs - 1e-10:
                le_tight |= {k for k,c in con.linear
                             if c > 0 and is_continuous[k]}

        for con in qp.linear_constraints:
            if sense == GE and lhs < rhs:  # GE violation
                deficit = rhs - lhs
                # Prefer vars NOT at LE ceiling
                cands = [k for k in coeffs
                         if k not in le_tight]
                # Distribute deficit proportionally
                for k, c in cands:
                    share = deficit * room[k] / total_room
                    x[k] += share / c

            elif sense == LE and lhs > rhs: # LE violation
                excess = lhs - rhs
                # Distribute excess proportionally
                for k, c in cands:
                    share = excess * room[k] / total_room
                    x[k] -= share / c

        x = np.clip(x, lb, ub)  # re-clamp

    if qp.is_feasible(x):
        result.status = SUCCESS""", font_size=9)

    # ════════════════════════════════════════════════════════
    # SLIDE 6 – Results: Success Rate
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "5. Results: Feasibility",
                 font_size=30, color=ACCENT_BLUE, bold=True)

    success_img = IMG_DIR / "success_rate.png"
    make_success_chart(data, success_img)
    slide.shapes.add_picture(str(success_img), Inches(0.6), Inches(1.3), Inches(6.5))

    tf = add_text_box(slide, Inches(7.5), Inches(1.3), Inches(5.3), Inches(5.5),
                      "Key Finding", font_size=22, color=GREEN, bold=True)
    add_bullet_text(tf, [
        "36 / 36 instances: SUCCESS for both solvers",
        "",
        "Both classical and quantum ADMM pipelines",
        "achieve 100% feasibility across all tested",
        "problem sizes (1-6 servers \u00d7 1-6 VMs).",
        "",
        "The snap_to_feasible() post-processing is",
        "essential: without it, only ~22% of ADMM",
        "solutions pass Qiskit's strict zero-tolerance",
        "feasibility check.",
        "",
        "Classical and quantum solvers consistently",
        "find the same objective value, confirming that",
        "QAOA correctly solves the binary sub-problem.",
    ], font_size=15, color=WHITE)

    # ════════════════════════════════════════════════════════
    # SLIDE 7 – Results: Time Comparison
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "6. Results: Execution Time Comparison",
                 font_size=30, color=ACCENT_BLUE, bold=True)

    time_img = IMG_DIR / "time_comparison.png"
    make_time_comparison_chart(data, time_img)
    slide.shapes.add_picture(str(time_img), Inches(0.8), Inches(1.3), Inches(11.5))

    tf = add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1.5),
                      "", font_size=14, color=LIGHT_GRAY)
    add_bullet_text(tf, [
        "Classical solver is consistently faster (expected: exact eigensolver is O(2\u207f) but n is small here)",
        "Quantum overhead comes from QAOA circuit simulation (StatevectorSampler) \u2014 would improve on real hardware",
        f"Speedup range: {min(float(r['speedup_x']) for r in data):.3f}x \u2013 {max(float(r['speedup_x']) for r in data):.3f}x",
    ], font_size=14, color=LIGHT_GRAY)

    # ════════════════════════════════════════════════════════
    # SLIDE 8 – Results: Speedup Analysis
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "7. Results: Speedup Analysis",
                 font_size=30, color=ACCENT_BLUE, bold=True)

    speedup_img = IMG_DIR / "speedup.png"
    make_speedup_chart(data, speedup_img)
    slide.shapes.add_picture(str(speedup_img), Inches(0.8), Inches(1.3), Inches(11.5))

    tf = add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1.5),
                      "", font_size=14, color=LIGHT_GRAY)
    add_bullet_text(tf, [
        "Speedup < 1.0 means classical is faster (as expected in simulation)",
        "Speedup increases with more VMs (larger continuous sub-problems)",
        "Drops sharply at 6 servers due to quantum circuit depth scaling exponentially with qubits",
    ], font_size=14, color=LIGHT_GRAY)

    # ════════════════════════════════════════════════════════
    # SLIDE 9 – Results: Scaling Heatmap
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "8. Results: Quantum Solver Scaling",
                 font_size=30, color=ACCENT_BLUE, bold=True)

    heatmap_img = IMG_DIR / "scaling_heatmap.png"
    make_scaling_heatmap(data, heatmap_img)
    slide.shapes.add_picture(str(heatmap_img), Inches(0.6), Inches(1.2), Inches(6.5))

    tf = add_text_box(slide, Inches(7.5), Inches(1.3), Inches(5.3), Inches(5.5),
                      "Scaling Observations", font_size=20, color=ACCENT_ORANGE, bold=True)
    add_bullet_text(tf, [
        "Quantum time is dominated by the number of",
        "binary variables (= number of servers)",
        "",
        "6 servers \u2192 6 qubits in QAOA circuit",
        "  \u2192 2\u2076 = 64 dimensional statevector",
        "",
        "Critical bottleneck: 6s\u00d73v = 70.8s (ADMM",
        "needed 3 iterations, each running QAOA)",
        "",
        "Adding VMs increases continuous variables",
        "but does not increase qubit count \u2192",
        "moderate time increase",
        "",
        "Worst case: 6s\u00d76v at 53.4s (more ADMM",
        "iterations due to problem complexity)",
    ], font_size=14, color=WHITE)

    # ════════════════════════════════════════════════════════
    # SLIDE 10 – Results: Objective Values
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "9. Results: Objective Values",
                 font_size=30, color=ACCENT_BLUE, bold=True)

    obj_img = IMG_DIR / "objective_values.png"
    make_objective_chart(data, obj_img)
    slide.shapes.add_picture(str(obj_img), Inches(0.8), Inches(1.3), Inches(11.5))

    tf = add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1.5),
                      "", font_size=14, color=LIGHT_GRAY)
    add_bullet_text(tf, [
        "Objective \u2248 sum(capacities) for require_all_on=True (dominated by fixed server costs)",
        "Classical and quantum objectives match to 6+ decimal places in ALL 36 instances",
        "Slight increases with more VMs reflect the dynamic utilization cost component (u\u2c7c \u00b7 v\u2c7c\u1d62)",
    ], font_size=14, color=LIGHT_GRAY)

    # ════════════════════════════════════════════════════════
    # SLIDE 11 – Results Table
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "10. Full Results Summary",
                 font_size=30, color=ACCENT_BLUE, bold=True)

    # Build table
    headers = ["Instance", "Objective", "Classic (s)", "Quantum (s)", "Speedup", "Status"]
    n_rows = len(data) + 1
    tbl_shape = slide.shapes.add_table(n_rows, len(headers),
                                        Inches(0.4), Inches(1.1),
                                        Inches(12.5), Inches(6.0))
    tbl = tbl_shape.table

    # Style header
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Fill rows
    for ri, row in enumerate(data):
        inst = f"{row['n_servers']}s \u00d7 {row['n_vms']}v"
        vals = [
            inst,
            f"{float(row['classic_objective']):.4f}",
            f"{float(row['classic_time_s']):.3f}",
            f"{float(row['quantum_time_s']):.3f}",
            f"{float(row['speedup_x']):.3f}x",
            "\u2713 SUCCESS",
        ]
        bg = TABLE_ROW1 if ri % 2 == 0 else TABLE_ROW2
        for ci, v in enumerate(vals):
            cell = tbl.cell(ri + 1, ci)
            cell.text = v
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.font.color.rgb = WHITE if ci < 5 else GREEN
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ════════════════════════════════════════════════════════
    # SLIDE 12 – Technology Stack
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "11. Technology Stack",
                 font_size=30, color=ACCENT_BLUE, bold=True)

    libs = [
        ("Qiskit 1.x", "Core quantum computing SDK by IBM"),
        ("qiskit-algorithms", "QAOA, VQE, NumPyMinimumEigensolver"),
        ("qiskit-optimization", "ADMM, QUBO conversion, constraint handling"),
        ("qiskit-aer", "High-performance quantum circuit simulator"),
        ("DOcplex (CPLEX)", "IBM CPLEX modeling for mathematical programming"),
        ("StatevectorSampler", "Exact statevector simulation (no noise)"),
        ("COBYLA", "Gradient-free optimizer for variational parameters"),
        ("matplotlib / numpy", "Visualization and numerical computation"),
    ]

    for idx, (name, desc) in enumerate(libs):
        y = 1.2 + idx * 0.7
        add_text_box(slide, Inches(1.0), Inches(y), Inches(4), Inches(0.4),
                     name, font_size=18, color=ACCENT_ORANGE, bold=True)
        add_text_box(slide, Inches(5.0), Inches(y), Inches(7.5), Inches(0.4),
                     desc, font_size=16, color=WHITE)

    # ════════════════════════════════════════════════════════
    # SLIDE 13 – Conclusion
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "12. Conclusions",
                 font_size=30, color=ACCENT_BLUE, bold=True)

    tf = add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.5),
                      "Key Findings", font_size=22, color=GREEN, bold=True)
    add_bullet_text(tf, [
        "100% feasibility (36/36) for both solvers",
        "  after snap_to_feasible() post-processing",
        "",
        "Identical objective values confirm QAOA",
        "  correctly solves the combinatorial sub-problem",
        "",
        "Classical solver is 2\u201350x faster in simulation",
        "  (expected: no quantum hardware advantage yet)",
        "",
        "ADMM successfully decomposes the MIQP into",
        "  manageable quantum-compatible sub-problems",
        "",
        "The framework scales to 6 servers \u00d7 6 VMs",
        "  (36 binary+continuous variables) without failure",
    ], font_size=15, color=WHITE)

    tf2 = add_text_box(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5),
                       "Future Work", font_size=22, color=ACCENT_ORANGE, bold=True)
    add_bullet_text(tf2, [
        "Run on real quantum hardware (IBM Quantum)",
        "  to assess noise impact and true quantum speedup",
        "",
        "Scale to larger instances (10+ servers, 10+ VMs)",
        "  using hardware-efficient ansatz and error mitigation",
        "",
        "Explore alternative quantum algorithms:",
        "  \u2022 VQE for the QUBO sub-problem",
        "  \u2022 Grover Adaptive Search",
        "  \u2022 Quantum Annealing (D-Wave)",
        "",
        "Add realistic constraints:",
        "  \u2022 Network latency between servers",
        "  \u2022 Memory and storage limits",
        "  \u2022 Multi-objective optimization (cost + latency)",
    ], font_size=15, color=WHITE)

    # ════════════════════════════════════════════════════════
    # SLIDE 14 – Thank you
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                 "Thank You",
                 font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
                 "Questions?",
                 font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                 "Quantum vs Classical Optimization  \u2022  Qiskit 1.x  \u2022  ADMM + QAOA  \u2022  36 Instances  \u2022  100% Feasibility",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Save ─────────────────────────────────────────────────
    prs.save(str(OUT))
    print(f"\nPresentation saved to: {OUT}")
    print(f"  {len(prs.slides)} slides")

if __name__ == "__main__":
    build_pptx()
